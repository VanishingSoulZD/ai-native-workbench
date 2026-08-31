# -*- coding: utf-8 -*-
"""
Phase 8 Task 4 — Build Landscape Explorer (HTML) and Executive Presentation (PPTX).

This is a BUILD-TIME generator. It renders validated Phase 8 research assets into
presentation form. It introduces NO new research: every string below is transcribed
from 08-dataset/products.csv, 08-dataset/candidates.csv, 08-sources.md and
08-research-note.md. Phase 0-7 research files are never read or modified by this script.

Run:  python3 build_assets.py
Outputs:
  landscape.html
  executive-summary.pptx
"""
import json
import os

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_HTML = os.path.join(HERE, "landscape.html")
OUT_PPTX = os.path.join(HERE, "executive-summary.pptx")

AS_OF = "2026-08-31"
RESEARCH_CUTOFF = "August 2026"

# ---------------------------------------------------------------------------
# 1. Top 10 products (locked order, transcribed from products.csv)
#    capability status: Confirmed / Partial / Unknown / NotPrimary
# ---------------------------------------------------------------------------
CAP_ORDER = ["planning", "context", "tools", "execution", "verification", "repair",
             "long_running", "multi_agent", "memory", "mcp", "skills", "sandbox", "cloud_agent"]
CAP_LABEL = {
    "planning": "Planning", "context": "Context", "tools": "Tools",
    "execution": "Execution", "verification": "Verification", "repair": "Repair",
    "long_running": "Long-running", "multi_agent": "Multi-agent", "memory": "Memory",
    "mcp": "MCP", "skills": "Skills", "sandbox": "Sandbox", "cloud_agent": "Cloud Agent",
}

def P(**kw):
    kw.setdefault("capabilities", {})
    return kw

PRODUCTS = [
    P(id="PF-01", name="Claude Code", company="Anthropic", position=1,
      paradigm="Terminal SWE Agent",
      categories=["Terminal SWE Agent", "Open / Provider-agnostic Harness"],
      leadership=["Market Adoption Leader", "Terminal SWE Agent Leader"],
      evidenceGrade="A", confidence="High",
      stratum="Stratum A — independently observable mainstream adoption",
      surfaces=["CLI", "IDE", "Desktop", "Cloud", "Terminal", "Slack/CI"],
      surfacesRaw="CLI, VS Code, JetBrains, desktop, claude.ai/code, Remote Control, Slack/CI",
      blurb="Strongest comparable independent adoption signal (JetBrains 39% work adoption); made the terminal a persistent agent operating environment.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Confirmed","mcp":"Confirmed","skills":"Confirmed","sandbox":"Partial","cloud_agent":"Confirmed"}),
    P(id="PF-02", name="Codex", company="OpenAI", position=2,
      paradigm="Delegated Multi-agent SWE",
      categories=["Delegated Multi-agent SWE", "Terminal SWE Agent"],
      leadership=["Delegated Multi-agent Leader"],
      evidenceGrade="A", confidence="High",
      stratum="Stratum A — independently observable mainstream adoption",
      surfaces=["App", "CLI", "IDE", "Cloud", "Web"],
      surfacesRaw="Codex app, CLI, IDE extension, cloud/web/GitHub-connected workflows",
      blurb="Cross-surface agent system (app/CLI/IDE/cloud); treats agent orchestration, not code generation, as the product. >5M WAU is a VENDOR CLAIM.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Confirmed","skills":"Confirmed","sandbox":"Confirmed","cloud_agent":"Confirmed"}),
    P(id="PF-03", name="Cursor", company="Anysphere / SpaceX", position=3,
      paradigm="AI-native IDE / Distributed Workspace",
      categories=["AI-native IDE / Distributed Workspace"],
      leadership=["AI-native IDE / Distributed Workspace Leader"],
      evidenceGrade="A", confidence="High",
      stratum="Stratum A — independently observable mainstream adoption",
      surfaces=["IDE", "CLI", "Desktop", "Cloud", "Mobile", "Slack/CI", "Self-hosted"],
      surfacesRaw="Desktop IDE, CLI, web/agents, mobile, Slack/GitHub/Linear entry points, Cloud Agents, self-hosted workers",
      blurb="AI-native IDE extended into Cloud Agents, parallel/background execution and self-hosted workers. Aug-2026 ownership change is strategic context, NOT a capability regression.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Confirmed","skills":"Confirmed","sandbox":"Confirmed","cloud_agent":"Confirmed"}),
    P(id="PF-04", name="GitHub Copilot", company="GitHub / Microsoft", position=4,
      paradigm="GitHub Lifecycle Agent",
      categories=["GitHub Lifecycle Agent", "Delegated Multi-agent SWE"],
      leadership=["GitHub Lifecycle Leader"],
      evidenceGrade="A", confidence="High",
      stratum="Stratum A — independently observable mainstream adoption",
      surfaces=["IDE", "CLI", "Cloud", "Desktop", "Mobile", "Web"],
      surfacesRaw="VS Code, Visual Studio, JetBrains, Neovim, CLI, GitHub.com, desktop app, mobile/remote-control",
      blurb="Owns the software-delivery graph (issue→code→PR); workflow integration is the moat candidate. 4.7M paid subs / 77K+ orgs is a VENDOR CLAIM.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Confirmed","mcp":"Confirmed","skills":"Confirmed","sandbox":"Confirmed","cloud_agent":"Confirmed"}),
    P(id="PF-05", name="Devin", company="Cognition", position=5,
      paradigm="Delegated Multi-agent SWE",
      categories=["Delegated Multi-agent SWE", "Enterprise Autonomous SWE Control Plane"],
      leadership=["Agent-Fleet / Autonomous SWE Leader"],
      evidenceGrade="B", confidence="Medium-High",
      stratum="Stratum B — strategic-scale agent platform",
      surfaces=["Desktop", "Cloud", "IDE", "API/MCP"],
      surfacesRaw="Devin Desktop, cloud sessions, editor/IDE workflow, API/MCP integrations",
      blurb="Packages autonomous coding as a unit of organizational labor (session/fleet management). 1M+ users / 4K+ enterprises are VENDOR CLAIMS.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Confirmed","mcp":"Confirmed","skills":"Partial","sandbox":"Partial","cloud_agent":"Confirmed"}),
    P(id="PF-06", name="Google Antigravity", company="Google", position=6,
      paradigm="Delegated Multi-agent SWE",
      categories=["Delegated Multi-agent SWE"],
      leadership=["Agent Command Center Leader"],
      evidenceGrade="A", confidence="Medium-High",
      stratum="Not assigned to a market stratum; 6% independent survey signal + Google/Gemini distribution",
      surfaces=["Desktop", "CLI", "Cloud", "Browser", "SDK/API"],
      surfacesRaw="Standalone desktop, CLI, SDK/runtime, integrated browser/Chrome and cloud/enterprise integration",
      blurb="Replaces the IDE as the primary human interface with an agent operations console. Gemini CLI is merged into this family (no double-count).",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Partial","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Confirmed","skills":"Confirmed","sandbox":"Confirmed","cloud_agent":"Partial"}),
    P(id="PF-07", name="Replit Agent", company="Replit", position=7,
      paradigm="Idea-to-Production Agent",
      categories=["Idea-to-Production Agent"],
      leadership=["Idea-to-Production Leader"],
      evidenceGrade="A", confidence="Medium-High",
      stratum="Stratum B — strategic-scale agent platform",
      surfaces=["Web", "Cloud", "Mobile", "Runtime/Deploy"],
      surfacesRaw="Web/cloud workspace with design, code, runtime, database and deployment; mobile/web creation flows",
      blurb="Idea→Production: the agent begins before the repository and finishes after the code, inside one managed runtime/deploy system. 50M+ users is PLATFORM-level, not agent-only.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Partial","skills":"Partial","cloud_agent":"Confirmed"}),
    P(id="PF-08", name="OpenCode", company="OpenCode / SST ecosystem", position=8,
      paradigm="Open / Provider-agnostic Harness",
      categories=["Open / Provider-agnostic Harness", "Terminal SWE Agent"],
      leadership=["Open-source Agent Harness Leader"],
      evidenceGrade="A", confidence="High",
      stratum="Stratum C — open / architecture-significant challenger",
      surfaces=["CLI/TUI", "Desktop", "Terminal"],
      surfacesRaw="CLI/TUI, desktop; client/server architecture supports remote driving",
      blurb="Open-source MIT, provider-agnostic harness (Build/Plan agents). Proves the harness can be a user-owned layer independent of model vendors.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Partial","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Confirmed","skills":"Confirmed","sandbox":"Partial","cloud_agent":"NotPrimary"}),
    P(id="PF-09", name="Qoder", company="Alibaba Cloud / Qoder", position=9,
      paradigm="Persistent Task-centric Agent",
      categories=["Persistent Task-centric Agent"],
      leadership=["Persistent Task-centric Leader"],
      evidenceGrade="A", confidence="Medium-High",
      stratum="Stratum B — strategic-scale agent platform",
      surfaces=["IDE", "CLI", "Cloud", "SDK/API", "Mobile"],
      surfacesRaw="IDE, CLI, JetBrains/VS Code ecosystem, QoderWork/CN surfaces, Cloud Agents, SDK/API",
      blurb="Persistent task runtime: planning/execution/verification/self-correction loop retains context across long tasks. 6M+ users / 100K+ businesses are VENDOR CLAIMS.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Confirmed","mcp":"Confirmed","skills":"Confirmed","sandbox":"Confirmed","cloud_agent":"Confirmed"}),
    P(id="PF-10", name="Factory", company="Factory AI", position=10,
      paradigm="Enterprise Autonomous SWE Control Plane",
      categories=["Enterprise Autonomous SWE Control Plane"],
      leadership=["Enterprise Deploy-anywhere Runtime Leader"],
      evidenceGrade="C", confidence="Medium",
      stratum="Stratum B — strategic-scale agent platform",
      surfaces=["Web/Control plane", "IDE", "Terminal", "Slack", "Browser", "CI/VM/K8s/Air-gapped"],
      surfacesRaw="Web/control plane, IDE integrations, terminal, Slack, browser, CI/VM/Kubernetes/air-gapped",
      blurb="Enterprise autonomous SWE via deployable Droids across laptops/CI/VM/K8s/air-gapped. Selected at #10 by research judgment despite lowest composite — NOT a weakness.",
      capabilities={"planning":"Confirmed","context":"Confirmed","tools":"Confirmed","execution":"Confirmed",
                    "verification":"Confirmed","repair":"Confirmed","long_running":"Confirmed","multi_agent":"Confirmed",
                    "memory":"Partial","mcp":"Partial","skills":"Partial","sandbox":"Partial","cloud_agent":"Partial"}),
]

# ---------------------------------------------------------------------------
# 2. Claim ledger (C001-C036, transcribed from 08-sources.md)
# ---------------------------------------------------------------------------
CLAIMS = [
    {"id":"C001","type":"Market Evidence","claim":"Claude Code has the strongest comparable independent adoption signal among measured products in the JetBrains May-Jul 2026 survey, at 39% work adoption.","evidence":"JetBrains surveyed 15,000+ professional developers (May-Jul 2026); 90% use AI coding agents at work weekly, 68% daily. Work-adoption: Claude Code 39%, Copilot 21%, Codex 16%, Cursor 12%, OpenCode 7%, Antigravity 6%.","sourceId":"S001","source":"JetBrains - AI Coding Agents: Adoption Trends","date":"2026-08","grade":"A","confidence":"High","phase":"P2;P7","scope":"Claude Code; market-wide"},
    {"id":"C002","type":"Market Evidence","claim":"The JetBrains adoption figures are multi-select survey responses that do not sum to 100%, and they are not market shares.","evidence":"Published figures sum to 110%, only possible under multi-select response. They are survey adoption signals, not shares of a single denominator.","sourceId":"S001","source":"JetBrains - AI Coding Agents: Adoption Trends","date":"2026-08","grade":"A","confidence":"High","phase":"P2;P7","scope":"Market-wide"},
    {"id":"C003","type":"Market Evidence","claim":"Adoption figures from different JetBrains survey waves are not directly comparable.","evidence":"April 2026 (Jan 2026 data) reports Copilot ~29%, Cursor ~18%, Claude Code ~18%; May-Jul 2026 wave reports Copilot 21%, Cursor 12%, Claude Code 39%. Instrument, period and referent differ.","sourceId":"S002;S001","source":"JetBrains - two survey waves","date":"2026-04; 2026-08","grade":"B","confidence":"High","phase":"P1;P2","scope":"Market-wide"},
    {"id":"C004","type":"Market Evidence","claim":"The ~90% adoption figure changes referent between Phase 1 and Phase 2/7 and the two formulations must not be merged.","evidence":"P1 §2.5 reports ~90% using at least one AI TOOL; P2 §3.1 and P7 §2.3 report 90% using AI CODING AGENTS at work weekly. Same number, different referent.","sourceId":"S201;S001","source":"01-candidate-universe.md §2.5; JetBrains Aug 2026","date":"2026-08","grade":"B","confidence":"High","phase":"P1;P2;P7","scope":"Market-wide"},
    {"id":"C005","type":"Judgment","claim":"There is no defensible single global AI Coding Agent market-share table for August 2026 using the evidence reviewed here.","evidence":"Survey adoption, paid seats, weekly active users, platform users, enterprise customer counts and GitHub stars measure different objects. No public source provides one denominator across commercial, OSS and China-market products.","sourceId":"S001;S031;S050;S056;S074;S208","source":"JetBrains; OpenAI; GitHub Newsroom; Devin; Qoder; 07-decision.md §4.3","date":"2026-08-31","grade":"B","confidence":"High","phase":"P7","scope":"Market-wide"},
    {"id":"C006","type":"Judgment","claim":"AI Coding Agent is best understood as an umbrella market of agentic software-engineering systems, not a single homogeneous product category: the technical substrate is converging while the product boundary is diverging.","evidence":"All ten products share intent, context, reasoning, tool use, execution, verification and repair, but each manages a different primary work object (task, repository, issue/PR, agent session, product, enterprise job).","sourceId":"S207;S208","source":"06-cross-product-analysis.md; 07-decision.md §3","date":"2026-08-31","grade":"A","confidence":"High","phase":"P6;P7","scope":"Market-wide"},
    {"id":"C007","type":"Judgment","claim":"Model -> Agent System -> Workflow is the dominant strategic direction; durable differentiation is moving above the model layer.","evidence":"Harness-sensitive benchmark outcomes plus product architecture show the competitive unit is Model + Harness + Runtime + Context/Memory + Tools + Verification + Workflow, not model identity alone.","sourceId":"S014;S023;S085;S208","source":"Terminal-Bench 2.1; Claude Code docs; Factory data-flow; 07-decision.md §3","date":"2026-08-31","grade":"B","confidence":"Medium-High","phase":"P5;P6;P7","scope":"Market-wide"},
    {"id":"C008","type":"Judgment","claim":"The agent harness is already a first-class strategic layer, and observed agent outcomes are scaffold/harness sensitive.","evidence":"P4 separates model, harness, tools, context and runtime across all ten products. P5 shows the same model producing materially different results under minimal-bash, SWE-agent-style, Claude Code and Codex CLI harnesses.","sourceId":"S023;S029;S085;S206","source":"Claude Code docs; Codex upgrades; Factory data-flow; 05-benchmarks.md §16","date":"2026-08-31","grade":"A","confidence":"High","phase":"P4;P5","scope":"Market-wide"},
    {"id":"C009","type":"Judgment","claim":"Runtime and sandbox are becoming part of the product itself, because runtime constraints define what autonomy is actually possible.","evidence":"Cursor Cloud Agents treat isolated VMs, dependencies, secrets, network, browser/desktop access and artifacts as prerequisites. Factory documents Droid execution across laptop, CI, VM, Kubernetes and air-gapped. Qoder Cloud Agents run in persistent cloud containers.","sourceId":"S040;S037;S084;S077","source":"Cursor Cloud Agents; Cursor self-hosted; Factory deployment patterns; Qoder Cloud Agent docs","date":"2026","grade":"A","confidence":"High","phase":"P4;P6","scope":"Cursor; Factory; Qoder"},
    {"id":"C010","type":"Judgment","claim":"Workflow integration is one of the strongest candidates for a durable long-term moat, but the duration and strength of that moat remain unproven.","evidence":"Workflow integration combines technical integration, switching costs, context, distribution and organizational process ownership. P7 records it as a strong candidate, not an established economic law.","sourceId":"S045;S064;S084;S208","source":"GitHub cloud agent docs; Replit Agent 4; Factory deployment patterns; 07-decision.md §5.7","date":"2026-08-31","grade":"B","confidence":"Medium-High","phase":"P7","scope":"Market-wide"},
    {"id":"C011","type":"Analysis","claim":"Basic MCP and tool-protocol support is moving toward commodity status; differentiation shifts to tool quality, permissioning, reliability and workflow integration.","evidence":"MCP support is Confirmed across all ten Top 10 products. Protocol support alone therefore no longer discriminates between them.","sourceId":"S023;S040;S069;S078;S208","source":"Claude Code docs; Cursor docs; OpenCode repo; Qoder slash reference; 07-decision.md §5.5","date":"2026-08-31","grade":"B","confidence":"High","phase":"P4;P7","scope":"Top 10"},
    {"id":"C012","type":"Analysis","claim":"The market is likely to commoditize individual agent primitives faster than complete agent systems; durable differentiation comes from composition.","evidence":"P7 classifies 8 capabilities as Commodity, 6 as Differentiating and 6 as Potential moat. Basic code generation, repository search, terminal access, basic planning, multi-file editing, MCP support and skills are all commodity or commoditizing.","sourceId":"S208","source":"07-decision.md §9","date":"2026-08-31","grade":"B","confidence":"High","phase":"P7","scope":"Market-wide"},
    {"id":"C013","type":"Benchmark Evidence","claim":"Automated SWE-bench grading can materially overstate real maintainer acceptance; automated grader pass rates averaged about 24.2 percentage points above maintainer merge decisions in the METR sample.","evidence":"METR reviewed 296 AI-generated PRs from 3 SWE-bench Verified repositories vs 47 human merged PRs. About half of automated-pass PRs would not be merged directly.","sourceId":"S006","source":"METR - Many SWE-bench-Passing PRs Would Not Be Merged into Main","date":"2026-03-10","grade":"A","confidence":"High","phase":"P5","scope":"Benchmark validity"},
    {"id":"C014","type":"Benchmark Evidence","claim":"Benchmark results do not establish developer productivity, and no clean unbiased estimate of current productivity uplift is available.","evidence":"METR's 2025 RCT found experienced OSS developers ~19% SLOWER under early-2025 AI conditions. METR's 2026-02 methodology update reports severe selection effects that prevent reading later-study data as an unbiased current uplift estimate.","sourceId":"S007;S008","source":"METR - 2025 RCT; METR - 2026-02 uplift update","date":"2025-07-10; 2026-02-24","grade":"A","confidence":"High","phase":"P5","scope":"Benchmark validity"},
    {"id":"C015","type":"Benchmark Evidence","claim":"The METR transcript analysis indicates a task time-savings factor of 1.5x-13x, but METR itself labels this a soft upper bound.","evidence":"5,305 Claude Code transcripts from 7 technical workers (Jan 2026). Caveats: task substitution, task selection effects, workers only use AI where helpful, saved time is not equivalent value.","sourceId":"S009","source":"METR - Analyzing coding agent transcripts","date":"2026-02-17","grade":"B","confidence":"Medium-High","phase":"P5","scope":"Claude Code"},
    {"id":"C016","type":"Judgment","claim":"No single public benchmark currently measures the full value of an AI Coding Agent product.","evidence":"The P5 capability coverage matrix shows Product UX, team collaboration, enterprise ROI/TCO and long-term memory quality are not measured by any reviewed benchmark; economic realism and maintainer review are covered by at most one or two.","sourceId":"S012;S013;S014;S015;S016;S017;S018;S019;S020;S206","source":"SWE-bench; ...; 05-benchmarks.md §15","date":"2026-08-31","grade":"A","confidence":"High","phase":"P5","scope":"Benchmark validity"},
    {"id":"C017","type":"Benchmark Evidence","claim":"The Kotlin Benchmark first public iteration reports Claude Code + Opus 4.7 xhigh at 85.71% (90/105), Junie + Opus 4.7 max at 81.9% and Codex + GPT-5.5 xhigh at 81.9%; these are AGENT + MODEL configurations, not product results.","evidence":"JetBrains published resolution rates as agent/model setups on 105 containerized Kotlin tasks. P5 states these are a first public benchmark run, not a permanent product ranking.","sourceId":"S018","source":"JetBrains - Kotlin Benchmark for AI Coding Agents","date":"2026-07","grade":"B","confidence":"Medium-High","phase":"P5","scope":"Benchmark results"},
    {"id":"C018","type":"Benchmark Evidence","claim":"Terminal-Bench 2.1 shows the same benchmark family producing materially different results across agent+model configurations.","evidence":"May 2026 2.1: GPT-5.3-Codex + Codex CLI 73.3%->79.1%; Opus 4.6 + Claude Code 58.0%->70.1%; Gemini 3.1 Pro + Terminus 2 63.0%->70.7%. Results must be pinned to benchmark version and submission date.","sourceId":"S014","source":"Terminal-Bench 2.1 release note and leaderboard","date":"2026-05-06","grade":"B","confidence":"Medium-High","phase":"P5","scope":"Benchmark results"},
    {"id":"C019","type":"Judgment","claim":"Phase 3 Top 10 is locked and is a Market Leaders / Representative Leaders research ranking, not a user-count Top 10 and not a capability ranking.","evidence":"Charter §3.2 defines the final research population as Market Significance x Technology/Product Significance. P5 §22 and P6 confirm the Top 10 unchanged.","sourceId":"S200;S203;S204;S206","source":"Charter §3.2; 03-ranking-methodology.md; 03-top10-selection.md; 05-benchmarks.md","date":"2026-08-31","grade":"A","confidence":"High","phase":"P3;P5;P6","scope":"Top 10"},
    {"id":"C020","type":"Judgment","claim":"Factory is selected for the locked Top 10 despite the lowest composite score among the ten, on representative workflow value and enterprise strategic significance.","evidence":"P3 §8.1 records the divergence: the Droids / autonomous enterprise SWE workflow is an independent paradigm; enterprise positioning does not fully overlap with Copilot or Devin; exclusion would over-concentrate the Top 10 on consumer/IDE/CLI products.","sourceId":"S204;S087","source":"03-top10-selection.md §8.1; Factory Series C","date":"2026-08-31","grade":"C","confidence":"Medium","phase":"P3","scope":"Factory"},
    {"id":"C021","type":"Judgment","claim":"Qoder's selection is evidence-backed rather than the result of a geographic quota.","evidence":"P3 §8.5 grounds selection in IDE + CLI + Cloud Agent surface coverage, memory/skills/MCP/browser/batch primitives, the Tongyi Lingma to Qoder lineage and the AI IDE to agent platform evolution.","sourceId":"S204;S074;S082","source":"03-top10-selection.md §8.5; Qoder changelog; Qoder CN billing docs","date":"2026-08-31","grade":"B","confidence":"Medium-High","phase":"P3","scope":"Qoder"},
    {"id":"C022","type":"Product Evidence","claim":"Amazon Q Developer CLI has been rebranded to Kiro and Amazon Q Developer IDE plugin support ends on 2027-04-30, so Q Developer is not a standalone modern product family for ranking purposes.","evidence":"AWS documentation states the CLI rebrand and the IDE plugin end-of-support date, and directs users to upgrade to Kiro.","sourceId":"S089;S090","source":"AWS - Upgrade to Kiro; AWS - Q Developer IDE end of support","date":"2026","grade":"A","confidence":"High","phase":"P2;P3","scope":"Amazon Q Developer; Kiro"},
    {"id":"C023","type":"Product Evidence","claim":"Gemini CLI has migrated to Antigravity CLI and shares the same harness and agent with Antigravity 2.0, so it must not be double-counted as a separate Google product family.","evidence":"Google's 2026-05-19 developer blog describes the transition, including transfer of the Gemini CLI user community, stars and contributors. I/O 26 materials state the CLI provides the same harness and agent as Antigravity 2.0.","sourceId":"S062;S061","source":"Google Developers Blog; Google Cloud - I/O 26 agent developer news","date":"2026-05","grade":"A","confidence":"High","phase":"P2;P3","scope":"Google Antigravity"},
    {"id":"C024","type":"Product Evidence","claim":"Devin Desktop is the new name for Windsurf, so Windsurf is not a separate market candidate.","evidence":"Cognition's official Devin Desktop pages state the renaming directly.","sourceId":"S052","source":"Devin Desktop","date":"2026","grade":"A","confidence":"High","phase":"P1;P4","scope":"Devin"},
    {"id":"C025","type":"Product Evidence","claim":"Qoder CN is the renamed continuation of the former Tongyi Lingma product line, renamed on 2026-05-20, so Tongyi Lingma is not separately counted.","evidence":"Qoder CN billing documentation records the rename; Alibaba/Qoder product materials describe the 2026 naming transition.","sourceId":"S082","source":"Qoder CN - billing description","date":"2026-05-20","grade":"A","confidence":"High","phase":"P1;P4","scope":"Qoder"},
    {"id":"C026","type":"Product Evidence","claim":"Claude Code executes across three documented environments: local machine, Anthropic-managed cloud VMs, and Remote Control from a browser.","evidence":"Claude Code documentation describes the three execution environments, meaning the product is no longer CLI-only at the operating-model level.","sourceId":"S023","source":"Claude Code Docs - How Claude Code works","date":"2026","grade":"A","confidence":"High","phase":"P4","scope":"Claude Code"},
    {"id":"C027","type":"Product Evidence","claim":"Factory Droids can run across laptops, CI pipelines, VMs, Kubernetes and air-gapped environments, with devcontainers and VMs providing isolation.","evidence":"Factory's deployment-pattern documentation enumerates these targets; public guidance states higher-autonomy agents should run in sandboxed environments.","sourceId":"S084","source":"Factory - Deployment patterns","date":"2026","grade":"A","confidence":"High","phase":"P4;P6","scope":"Factory"},
    {"id":"C028","type":"Product Evidence","claim":"OpenCode is an open-source MIT agent harness that is explicitly provider-agnostic, with Build and Plan as distinct primary agents.","evidence":"The OpenCode repository and agents documentation describe the MIT licence, TUI focus, client/server architecture, Build/Plan primary agents and support for many providers plus local models.","sourceId":"S069;S070;S072","source":"OpenCode GitHub repository; agents docs; LLM package docs","date":"2026","grade":"A","confidence":"High","phase":"P4","scope":"OpenCode"},
    {"id":"C029","type":"Product Evidence","claim":"GitHub Copilot's Pro plan makes the third-party agents Claude Code and Codex available; this is a distribution fact, not a ranking input or an endorsement.","evidence":"Recorded in the P4 product-04 evidence ledger from the GitHub plans page. Not surfaced in P6 or P7.","sourceId":"S044","source":"GitHub Copilot plans","date":"2026","grade":"A","confidence":"High","phase":"P4","scope":"GitHub Copilot"},
    {"id":"C030","type":"Market Evidence","claim":"Codex reported more than 5 million weekly active users as of 2026-06-02; this is a vendor claim, not an independently normalized market denominator.","evidence":"OpenAI's own productivity report. Separate from the 16% independent survey signal.","sourceId":"S031","source":"OpenAI - Codex is becoming a productivity tool for everyone","date":"2026-06-02","grade":"B","confidence":"High","phase":"P2;P3","scope":"Codex"},
    {"id":"C031","type":"Market Evidence","claim":"Devin's 1M+ users and 4,000+ enterprise customers, Qoder's 6M+ users and 100K+ businesses, and Replit's 50M+ platform users are all vendor claims and must not be converted into independent market facts.","evidence":"Devin figures from vendor product pages; Qoder figures from the 2026-08-26 changelog; Replit's figure is explicitly platform-level rather than agent-only.","sourceId":"S056;S074;S066","source":"Devin product claims; Qoder changelog; Replit Pricing","date":"2026","grade":"C","confidence":"Medium-High","phase":"P2;P3","scope":"Devin; Qoder; Replit Agent"},
    {"id":"C032","type":"Analysis","claim":"Cursor's August 2026 corporate transition is strategic context and is not evidence of a product capability regression.","evidence":"Acquisition by SpaceX completed 2026-08-14; OpenAI announced model-supply wind-down 2026-08-28 with proposed 2026-11-12 cutoff; Reuters reported 2026-08-29. P3, P6, P7 treat this as ecosystem/model-supply risk rather than capability change.","sourceId":"S042;S005;S208","source":"Cursor - Joining SpaceX; Reuters 2026-08-29; 07-decision.md §2.3","date":"2026-08-29","grade":"A","confidence":"High","phase":"P4;P6;P7","scope":"Cursor"},
    {"id":"C033","type":"Judgment","claim":"The software-engineering work unit has clearly risen above the line/file level, but it has not universally reached whole-workflow autonomy.","evidence":"Product design shows delegated agent workstreams and parallel portfolios; independent productivity evidence remains incomplete and noisy. P7 separates 'already happened', 'emerging now' and 'not yet proven'.","sourceId":"S001;S008;S208","source":"JetBrains Aug 2026; METR uplift update; 07-decision.md §10","date":"2026-08-31","grade":"B","confidence":"High","phase":"P7","scope":"Market-wide"},
    {"id":"C034","type":"Analysis","claim":"The estimate that roughly 47% of produced code was fully agent-generated on average is a DERIVED midpoint calculation across survey buckets, not a measured fact.","evidence":"P7 §10.2 attributes the figure to a JetBrains companion analysis using midpoint calculations across buckets, with about 22% of developers in the >80% agent-generated group.","sourceId":"S123","source":"JetBrains - agent-generated code companion analysis","date":"2026-08","grade":"C","confidence":"Medium","phase":"P7","scope":"Market-wide"},
    {"id":"C035","type":"Unknown","claim":"End-to-end autonomous software engineering is NOT proven, and reliable low-supervision autonomous delivery remains unproven.","evidence":"METR's maintainer study shows automated pass is not maintainer acceptance; METR's methodology update blocks an unbiased productivity estimate. P7 §18 lists this under 'Not Yet Proven'.","sourceId":"S006;S008;S208","source":"METR maintainer study; METR uplift update; 07-decision.md §18","date":"2026-08-31","grade":"B","confidence":"High","phase":"P5;P7","scope":"Market-wide"},
    {"id":"C036","type":"Unknown","claim":"True agent-level market share cannot be established from the public evidence reviewed in this Case.","evidence":"No unified public denominator covers commercial, open-source and China-market coding agents. P7 §14 records this as a standing Unknown.","sourceId":"S001;S208","source":"JetBrains Aug 2026; 07-decision.md §14","date":"2026-08-31","grade":"C","confidence":"High","phase":"P2;P7","scope":"Market-wide"},
]

# ---------------------------------------------------------------------------
# 3. Structural content (transcribed from 08-research-note.md / Phases 6-7)
# ---------------------------------------------------------------------------
CATEGORIES = [
    {"name":"Terminal SWE Agent","desc":"The terminal/executable environment is where the agent observes, acts, verifies, repairs. The terminal becomes a persistent agent operating environment.","members":["Claude Code","Codex","OpenCode"]},
    {"name":"AI-native IDE / Distributed Workspace","desc":"The IDE stays a first-class human surface while execution moves to remote/background workers, isolated VMs and self-hosted infrastructure.","members":["Cursor"]},
    {"name":"GitHub Lifecycle Agent","desc":"The agent attaches to the software-delivery graph — repository, issue, PR, review, CI, enterprise policy — rather than a new environment.","members":["GitHub Copilot"]},
    {"name":"Delegated Multi-agent SWE","desc":"The task/workstream/session is the unit; the human supervises asynchronous executions, parallel agents and long-running work.","members":["Codex","GitHub Copilot","Devin","Google Antigravity"]},
    {"name":"Persistent Task-centric Agent","desc":"The task persists across specification, memory, execution, verification and cloud runtime rather than disappearing with one session.","members":["Qoder"]},
    {"name":"Idea-to-Production Agent","desc":"The agent begins before the repository and finishes after the code, inside one managed runtime/deployment system; the outcome is a running product.","members":["Replit Agent"]},
    {"name":"Enterprise Autonomous SWE Control Plane","desc":"The organization treats the agent as an operational unit running in approved infrastructure under policy and human/automated gates.","members":["Devin","Factory"]},
    {"name":"Open / Provider-agnostic Harness","desc":"The model provider is a replaceable substrate while the harness stays user-owned/open — an architecture alternative to model-vendor-owned stacks.","members":["Claude Code","OpenCode"]},
]

PARADIGMS = [
    ("Terminal SWE Agent","Claude Code, OpenCode (Codex partially). The terminal/executable environment is where the agent observes, acts, verifies, repairs."),
    ("AI-native IDE / Distributed Workspace","Cursor (Qoder overlaps). The IDE stays a first-class human surface while execution moves to remote/background workers."),
    ("GitHub Lifecycle Agent","GitHub Copilot. The agent attaches to the software-delivery graph rather than a new environment."),
    ("Delegated Multi-agent SWE","Codex, Devin, Antigravity. The task/workstream/session is the unit; the human supervises asynchronous executions."),
    ("Persistent Task-centric Agent","Qoder. The task persists across specification, memory, execution, verification and cloud runtime."),
    ("Idea-to-Production Agent","Replit Agent. The agent begins before the repository and finishes after the code, inside one managed runtime/deploy system."),
    ("Enterprise Autonomous SWE","Factory (Devin partially). The organization treats the agent as an operational unit running in approved infrastructure under policy."),
    ("Open / Provider-agnostic Harness","OpenCode. The model provider is a replaceable substrate while the harness stays user-owned/open."),
]

STRATEGIC_LAYERS = [
    {"name":"Model","importance":"Very high","differentiation":"Very high today","commoditization":"High–Medium","moat":"Medium unless proprietary access/cost/performance is durable","confidence":"High"},
    {"name":"Harness","importance":"Very high","differentiation":"High","commoditization":"Medium","moat":"High","confidence":"High"},
    {"name":"Runtime","importance":"Very high for autonomous agents","differentiation":"High","commoditization":"Medium","moat":"High (esp. enterprise)","confidence":"High"},
    {"name":"Context / Memory","importance":"Very high","differentiation":"Medium–High","commoditization":"Medium","moat":"Medium–High","confidence":"High"},
    {"name":"Tools / MCP","importance":"High as substrate","differentiation":"Medium","commoditization":"High","moat":"Low–Medium individually","confidence":"High"},
    {"name":"Orchestration","importance":"High and rising","differentiation":"High","commoditization":"Medium","moat":"High","confidence":"High"},
    {"name":"Workflow","importance":"Very high","differentiation":"Very high","commoditization":"Low–Medium","moat":"Very high","confidence":"High"},
    {"name":"Distribution","importance":"Very high","differentiation":"Very high","commoditization":"Low","moat":"Very high","confidence":"High"},
]

COMMODITIZATION = {
    "Commodity / rapidly commoditizing": ["Basic code generation","repository search / basic codebase context","terminal access","basic planning","multi-file editing","basic testing / command execution","MCP support (commoditizing)","Skills / reusable instructions (commoditizing)"],
    "Differentiating": ["Reliable verification / repair","long-running execution","context / memory quality","agent orchestration","environment integration","human steering / review UX"],
    "Potential moat": ["Runtime infrastructure","workflow integration (strong)","orchestration / control plane (strong)","organizational memory / context graph","verification / evaluation system","distribution (strong)"],
}

COMPETITIVE = [
    ("Direct", ["Claude Code ↔ Codex","Cursor ↔ Qoder","Devin ↔ Factory"]),
    ("Adjacent", ["Cursor ↔ Antigravity","Codex ↔ Devin","Copilot ↔ standalone agent platforms"]),
    ("Workflow alternatives", ["Replit ↔ repo-centric agents","Antigravity ↔ Codex / Devin","Copilot ↔ standalone agents"]),
    ("Architecture alternatives", ["OpenCode ↔ vendor-owned harnesses","Factory ↔ SaaS-only cloud agents","Cursor self-hosted ↔ vendor-hosted agents"]),
]

WORKFLOW = ["Code Completion","Code / File Editing","Issue / Task Resolution","Repository-level Execution",
            "Project / Product Work","Delegated Agent Workstream","Parallel Agent Portfolio","Engineering Workflow Automation"]

SCENARIOS = [
    ("A — Model Dominates","Frontier-model progress remains so rapid that model quality is the main differentiator and harness/runtime standardizes. Plausible, but no longer sufficient as the default market model.","Limiting evidence: multi-model routing, provider-agnostic products, harness-sensitive benchmarks, runtime-heavy systems."),
    ("B — Agent System Dominates","Harness + runtime + context + memory + tools + orchestration become the main differentiators around still-important models. Most supported current trajectory.","Confidence: Medium-High."),
    ("C — Workflow Platform Dominates","Agent systems embed in complete dev workflows (intake, spec, execution, review, deployment, governance, workforce management). Strong strategic direction, but not yet proven as the final winning structure.","Confidence: plausible, not proven."),
]
SCENARIO_ORDER = "Support ordering: Scenario B > Scenario C > Scenario A. This is a strategic judgment, NOT a calibrated probability forecast."

RISKS = [
    ("Market denominator limitations","No unified public denominator covers commercial, OSS and China-market coding agents; no defensible single global market-share table for Aug 2026 (C005, C036)."),
    ("Vendor-claim limitations","User/customer counts (Codex >5M WAU, Copilot 4.7M/77K, Devin 1M+/4K, Qoder 6M+/100K, Replit 50M platform, Factory scale) are vendor claims, not independent fact (C030, C031)."),
    ("Benchmark limitations","Automated grader pass rates overstate maintainer acceptance by ~24.2 pts; benchmarks do not establish productivity; no clean uplift estimate exists (C013, C014, C016)."),
    ("Model / provider dependency","Frontier improvements can quickly shift product capability; distribution can increase dependency on a model/provider ecosystem."),
    ("Corporate / product uncertainty","Families actively acquire, rename and expand surfaces (Q→Kiro, Gemini CLI→Antigravity, Windsurf→Devin Desktop, Tongyi Lingma→Qoder); Cursor ownership change is strategic context, not regression (C022–C025, C032)."),
    ("Hidden internal implementation","Planner architectures, routing policies, evaluator internals and success classifiers are not public for most vendors."),
    ("Future moat uncertainty","Workflow lock-in durability, long-horizon autonomy reliability, enterprise ROI and security outcomes remain under-measured or unproven (C035)."),
]

EXEC_JUDGMENT = ("AI Coding Agent is an umbrella market of agentic software-engineering systems. "
    "The technical substrate is converging while the product boundary is diverging. "
    "Competition is moving Model → Agent System → Workflow: durable differentiation increasingly lives in "
    "harness, runtime, context/memory, tools, orchestration, verification and workflow integration — not model identity alone.")

MARKET_STRUCTURE = [
    ("What the market is","AI Coding Agent is best read as the middle label (agent) expanding into Agentic Software Engineering Platform. The market is an umbrella of related systems, not one homogeneous category."),
    ("Market strata","Stratum A — independently observable mainstream adoption (Claude Code, Copilot, Codex, Cursor). Stratum B — strategic-scale agent platforms (Devin, Qoder, Replit, Factory). Stratum C — open / architecture-significant challengers (OpenCode). Strata are NOT quality tiers."),
    ("Market denominator","There is NO defensible single global AI Coding Agent market-share table for August 2026. Survey adoption, paid seats, WAU, platform users and enterprise counts measure different objects and must not be collapsed into false precision (JetBrains figures are multi-select, sum to 110%, and are NOT shares of a single denominator)."),
]

# ---------------------------------------------------------------------------
# 4. HTML rendering
# ---------------------------------------------------------------------------
DATA = {
    "asOf": AS_OF, "cutoff": RESEARCH_CUTOFF,
    "products": PRODUCTS, "claims": CLAIMS,
    "categories": CATEGORIES, "paradigms": [{"name":n,"desc":d} for n,d in PARADIGMS],
    "strategicLayers": STRATEGIC_LAYERS, "commoditization": COMMODITIZATION,
    "competitive": [{"group":g,"pairs":p} for g,p in COMPETITIVE],
    "workflow": WORKFLOW, "scenarios": [{"name":n,"desc":d,"note":nt} for n,d,nt in SCENARIOS],
    "scenarioOrder": SCENARIO_ORDER, "risks": [{"name":n,"desc":d} for n,d in RISKS],
    "capOrder": CAP_ORDER, "capLabel": CAP_LABEL,
    "execJudgment": EXEC_JUDGMENT, "marketStructure": [{"k":k,"v":v} for k,v in MARKET_STRUCTURE],
}

HTML_HEAD = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>2026 AI Coding Agent Landscape Explorer</title>
<style>
:root{
  --bg:#0f1115; --panel:#161a21; --panel2:#1d222b; --ink:#e8edf4; --muted:#97a3b6;
  --line:#2a313d; --accent:#5b9dff; --accent2:#7ee0c0; --warn:#ffb454; --bad:#ff6b6b; --good:#7ee0c0;
  --c:#3fb950; --p:#d29922; --u:#8b949e; --np:#586069;
}
*{box-sizing:border-box}
body{margin:0;font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,Helvetica,Arial,"PingFang SC","Microsoft YaHei",sans-serif;background:var(--bg);color:var(--ink);line-height:1.55}
a{color:var(--accent);text-decoration:none}
header.hero{padding:64px 24px 48px;text-align:center;background:radial-gradient(1200px 500px at 50% -10%,#1b2740,transparent),var(--bg);border-bottom:1px solid var(--line)}
header.hero h1{font-size:clamp(28px,4.5vw,46px);margin:0 0 12px;letter-spacing:.5px}
header.hero .sub{color:var(--muted);font-size:clamp(14px,2vw,18px);max-width:900px;margin:0 auto}
.badge{display:inline-block;margin-top:18px;padding:6px 12px;border:1px solid var(--line);border-radius:999px;color:var(--muted);font-size:12px}
nav.toc{position:sticky;top:0;z-index:30;background:rgba(15,17,21,.92);backdrop-filter:blur(8px);border-bottom:1px solid var(--line);display:flex;flex-wrap:wrap;gap:4px;padding:8px 12px;justify-content:center}
nav.toc a{font-size:12.5px;color:var(--muted);padding:5px 9px;border-radius:7px}
nav.toc a:hover{background:var(--panel2);color:var(--ink)}
main{max-width:1120px;margin:0 auto;padding:0 18px 80px}
section{padding:42px 0;border-bottom:1px solid var(--line)}
section h2{font-size:24px;margin:0 0 6px;display:flex;align-items:center;gap:10px}
section h2 .n{font-size:13px;color:var(--accent);border:1px solid var(--accent);border-radius:6px;padding:1px 7px}
section .lead{color:var(--muted);margin:0 0 22px;max-width:860px}
.lock{display:inline-block;font-size:11.5px;color:var(--warn);border:1px solid var(--warn);border-radius:6px;padding:2px 8px;margin-left:8px;vertical-align:middle}
.card{background:var(--panel);border:1px solid var(--line);border-radius:12px;padding:16px 18px;margin-bottom:14px}
.grid{display:grid;gap:14px}
.grid.cols2{grid-template-columns:repeat(2,1fr)}
.grid.cols3{grid-template-columns:repeat(3,1fr)}
@media(max-width:820px){.grid.cols2,.grid.cols3{grid-template-columns:1fr}}
.tag{display:inline-block;font-size:11.5px;padding:2px 8px;border-radius:999px;border:1px solid var(--line);color:var(--muted);margin:2px 4px 2px 0}
.tag.cat{background:#15233a;color:#bcd4ff}
.tag.lead{background:#1a2a22;color:#a8e6cf}
.tag.surf{background:#241f15;color:#f3d9a6}
.pos{font-weight:700;color:var(--accent)}
.rankrow{display:flex;gap:14px;align-items:flex-start}
.ranknum{flex:0 0 38px;height:38px;border-radius:9px;background:var(--panel2);border:1px solid var(--line);display:flex;align-items:center;justify-content:center;font-weight:700;color:var(--accent)}
table{border-collapse:collapse;width:100%;font-size:13.5px}
th,td{border:1px solid var(--line);padding:7px 9px;text-align:left;vertical-align:top}
th{background:var(--panel2);color:var(--ink)}
.cellC{background:rgba(63,185,80,.16);color:#9ff0b0}
.cellP{background:rgba(210,153,34,.18);color:#f0cf86}
.cellU{background:rgba(139,148,158,.12);color:#aeb6c2}
.cellN{background:rgba(88,96,105,.10);color:#7b8492}
.filters{position:sticky;top:46px;z-index:20;background:var(--panel2);border:1px solid var(--line);border-radius:12px;padding:12px 14px;margin-bottom:18px}
.filters .row{display:flex;flex-wrap:wrap;gap:8px;align-items:center;margin:6px 0}
.filters label{font-size:12px;color:var(--muted);min-width:96px}
.filters input[type=search]{flex:1;min-width:180px;background:var(--bg);border:1px solid var(--line);color:var(--ink);border-radius:8px;padding:7px 10px}
.chip{font-size:12px;padding:4px 10px;border-radius:999px;border:1px solid var(--line);background:var(--bg);color:var(--muted);cursor:pointer;user-select:none}
.chip.on{background:var(--accent);color:#06101f;border-color:var(--accent);font-weight:600}
select{background:var(--bg);color:var(--ink);border:1px solid var(--line);border-radius:8px;padding:6px 8px;font-size:12.5px}
.btn{font-size:12px;padding:4px 10px;border-radius:8px;border:1px solid var(--line);background:var(--bg);color:var(--muted);cursor:pointer}
.btn:hover{color:var(--ink)}
.lvlbar{height:8px;border-radius:5px;background:var(--panel2);overflow:hidden;margin-top:4px}
.hidden{display:none!important}
.modal-bg{position:fixed;inset:0;background:rgba(0,0,0,.6);display:flex;align-items:center;justify-content:center;z-index:50;padding:18px}
.modal{background:var(--panel);border:1px solid var(--line);border-radius:14px;max-width:720px;width:100%;max-height:86vh;overflow:auto;padding:20px 22px}
.modal h3{margin:0 0 12px;font-size:19px}
.kv{display:grid;grid-template-columns:140px 1fr;gap:6px 12px;font-size:13.5px;margin:10px 0}
.kv .k{color:var(--muted)}
.modal .close{float:right;cursor:pointer;color:var(--muted);font-size:22px;line-height:1}
.claimrow{display:flex;gap:10px;align-items:flex-start;padding:9px 0;border-bottom:1px solid var(--line);cursor:pointer}
.claimrow:hover{background:var(--panel2)}
.claimrow .cid{flex:0 0 54px;font-weight:700;color:var(--accent)}
.grade{font-size:11px;padding:1px 7px;border-radius:6px;border:1px solid var(--line)}
.gA{color:#9ff0b0;border-color:#2e6b3a}.gB{color:#bcd4ff;border-color:#2e4a6b}.gC{color:#f0cf86;border-color:#6b5a2e}.gD{color:#ffb0b0;border-color:#6b2e2e}
.note{font-size:12px;color:var(--muted);margin-top:8px}
.foot{text-align:center;color:var(--muted);font-size:12px;padding:30px 18px 60px}
.stage{display:flex;flex-wrap:wrap;align-items:center;gap:6px}
.stage .s{background:var(--panel2);border:1px solid var(--line);border-radius:8px;padding:7px 11px;font-size:13px}
.stage .arrow{color:var(--muted)}
.scenario{border-left:3px solid var(--accent);padding:10px 14px;background:var(--panel);border-radius:0 10px 10px 0;margin-bottom:12px}
.dot{display:inline-block;width:10px;height:10px;border-radius:50%;margin-right:6px;vertical-align:middle}
</style>
</head>
<body>
"""

HTML_BODY = """
<header class="hero">
  <h1>2026 AI Coding Agent Landscape Explorer</h1>
  <div class="sub">Market Structure · Product Paradigms · Agent Architecture · Workflow Evolution · Strategic Outlook</div>
  <div class="badge" id="asof"></div>
</header>

<nav class="toc">
  <a href="#exec">Executive</a>
  <a href="#market">Market</a>
  <a href="#top10">Top 10</a>
  <a href="#catmap">Category Map</a>
  <a href="#leadmap">Leadership</a>
  <a href="#layers">Strategic Layers</a>
  <a href="#commod">Commoditization</a>
  <a href="#comp">Competitive</a>
  <a href="#wf">Workflow</a>
  <a href="#scen">Scenarios</a>
  <a href="#risk">Risks/Unknowns</a>
  <a href="#sources">Evidence/Sources</a>
</nav>

<main>
  <section id="exec">
    <h2><span class="n">§</span>Executive Summary <span class="lock">Locked Phase 3 Top 10</span></h2>
    <p class="lead">A long-term knowledge asset rendered from the validated Phase 8 research record. This is a presentation layer, not a new research layer.</p>
    <div class="card" id="execJudgment"></div>
    <div class="grid cols3" id="execStats"></div>
  </section>

  <section id="market">
    <h2><span class="n">§</span>Market Structure</h2>
    <p class="lead">Umbrella market; substrate converging, product boundary diverging. No single global market-share table is defensible for Aug 2026.</p>
    <div id="marketStructure"></div>
  </section>

  <section id="top10">
    <h2><span class="n">§</span>Top 10 <span class="lock">Locked Phase 3 Top 10 — not a current global ranking</span></h2>
    <p class="lead">Market Significance × Technology/Product Significance. Locked order — do not re-sort, do not present as a new ranking. Use the filters to explore.</p>
    <div class="filters" id="filters"></div>
    <div id="top10list"></div>
  </section>

  <section id="catmap">
    <h2><span class="n">§</span>Category Map <span class="lock">Phase 6 categories — overlap allowed</span></h2>
    <p class="lead">Eight architecture paradigms. A product may belong to several; the label names its dominant architectural idea.</p>
    <div id="catmapView"></div>
  </section>

  <section id="leadmap">
    <h2><span class="n">§</span>Leadership Map <span class="lock">Category judgments — not a second overall ranking</span></h2>
    <p class="lead">Each category is led by its strongest representative under the reviewed evidence — not a universal "best overall".</p>
    <div id="leadmapView"></div>
  </section>

  <section id="layers">
    <h2><span class="n">§</span>Strategic Layer Model <span class="lock">8 ordered layers (Phase 7)</span></h2>
    <p class="lead">Model → Agent System → Workflow. Ordered, not ranked; no layer scores assigned.</p>
    <table id="layersTable"></table>
  </section>

  <section id="commod">
    <h2><span class="n">§</span>Capability Commoditization <span class="lock">3-level structure (Phase 7)</span></h2>
    <p class="lead">Commodity / rapidly commoditizing · Differentiating · Potential moat. The market commoditizes primitives faster than complete systems.</p>
    <div class="grid cols3" id="commodView"></div>
  </section>

  <section id="comp">
    <h2><span class="n">§</span>Competitive Structure <span class="lock">Relationships from Phase 6 only</span></h2>
    <p class="lead">Typed, not ranked. No new competitive relationships invented.</p>
    <div id="compView"></div>
  </section>

  <section id="wf">
    <h2><span class="n">§</span>Workflow Evolution</h2>
    <p class="lead">The work unit is rising; it has not universally reached whole-workflow autonomy.</p>
    <div class="stage" id="wfView"></div>
    <p class="note" id="wfNote"></p>
  </section>

  <section id="scen">
    <h2><span class="n">§</span>Strategic Scenarios <span class="lock">Phase 7 scenarios only</span></h2>
    <p class="lead">No new forecasts or market-share predictions.</p>
    <div id="scenView"></div>
    <p class="note" id="scenOrder"></p>
  </section>

  <section id="risk">
    <h2><span class="n">§</span>Risks / Unknowns</h2>
    <p class="lead">Evidence limits are carried forward, never resolved.</p>
    <div class="grid cols2" id="riskView"></div>
  </section>

  <section id="sources">
    <h2><span class="n">§</span>Evidence / Sources</h2>
    <p class="lead">Click any claim to inspect Claim · Evidence · Source · Date · Evidence Grade · Confidence · Phase. Source registry: 08-sources.md (36 claims, 131 sources).</p>
    <div id="claimsList"></div>
  </section>
</main>

<div class="modal-bg hidden" id="modal">
  <div class="modal" id="modalBody"></div>
</div>

<div class="foot" id="foot"></div>

<script>
const $=s=>document.querySelector(s), $$=s=>[...document.querySelectorAll(s)];
const capClass=s=>s==='Confirmed'?'cellC':s==='Partial'?'cellP':s==='NotPrimary'?'cellN':'cellU';
const capShort=s=>s==='Confirmed'?'C':s==='Partial'?'P':s==='NotPrimary'?'–':'U';
const nameToProduct={}; DATA.products.forEach(p=>nameToProduct[p.name]=p);

// ---- filters ----
const state={q:'',paradigm:new Set(),category:new Set(),surface:'',grade:new Set(),cap:new Set()};
function uniq(arr){return [...new Set(arr)].sort();}
const allParadigms=uniq(DATA.products.map(p=>p.paradigm));
const allCategories=uniq(DATA.products.flatMap(p=>p.categories));
const allSurfaces=uniq(DATA.products.flatMap(p=>p.surfaces));
const allLead=uniq(DATA.products.flatMap(p=>p.leadership));
const allGrades=uniq(DATA.products.map(p=>p.evidenceGrade));

function buildFilters(){
  const f=$('#filters');
  let html='<div class="row"><label>Search</label><input type="search" id="fq" placeholder="name or company…"></div>';
  html+='<div class="row"><label>Workflow paradigm</label>'+allParadigms.map(p=>`<span class="chip" data-k="paradigm" data-v="${p}">${p}</span>`).join('')+'</div>';
  html+='<div class="row"><label>Category</label>'+allCategories.map(c=>`<span class="chip" data-k="category" data-v="${c}">${c}</span>`).join('')+'</div>';
  html+='<div class="row"><label>Surface</label><select id="fsurf"><option value="">Any</option>'+allSurfaces.map(s=>`<option>${s}</option>`).join('')+'</select></div>';
  html+='<div class="row"><label>Leadership role</label>'+allLead.map(l=>`<span class="chip" data-k="lead" data-v="${l}">${l}</span>`).join('')+'</div>';
  html+='<div class="row"><label>Evidence grade</label>'+['A','B','C','D'].filter(g=>allGrades.includes(g)).map(g=>`<span class="chip" data-k="grade" data-v="${g}">${g}</span>`).join('')+' <span class="chip" data-k="grade" data-v="">Any</span></div>';
  html+='<div class="row"><label>Capability</label><span class="chip" data-k="cap" data-v="Confirmed">Confirmed</span><span class="chip" data-k="cap" data-v="Partial">Partial</span><span class="chip" data-k="cap" data-v="Unknown">Unknown</span> <span class="btn" id="freset">Reset</span></div>';
  f.innerHTML=html;
  $('#fq').addEventListener('input',e=>{state.q=e.target.value.toLowerCase();renderProducts();});
  $('#fsurf').addEventListener('change',e=>{state.surface=e.target.value;renderProducts();});
  $$('#filters .chip').forEach(c=>c.addEventListener('click',()=>{
    const k=c.dataset.k,v=c.dataset.v,set=state[k];
    if(v===''){set.clear();} else if(set.has(v))set.delete(v); else set.add(v);
    c.classList.toggle('on', v!=='' && set.has(v));
    if(v==='' ){} renderProducts();
  }));
  $('#freset').addEventListener('click',()=>{state.q='';state.paradigm.clear();state.category.clear();state.surface='';state.grade.clear();state.cap.clear();$('#fq').value='';$('#fsurf').value='';$$('#filters .chip').forEach(c=>c.classList.remove('on'));renderProducts();});
}

function passFilter(p){
  if(state.q && !(p.name.toLowerCase().includes(state.q)||p.company.toLowerCase().includes(state.q))) return false;
  if(state.paradigm.size && !state.paradigm.has(p.paradigm)) return false;
  if(state.category.size && !p.categories.some(c=>state.category.has(c))) return false;
  if(state.surface && !p.surfaces.includes(state.surface)) return false;
  if(state.grade.size && !state.grade.has(p.evidenceGrade)) return false;
  return true;
}
function filtered(){return DATA.products.filter(passFilter);}

function renderProducts(){
  const list=filtered();
  $('#top10list').innerHTML = list.length? list.map(p=>productCard(p)).join('')
    : '<div class="card">No product matches the current filters.</div>';
  renderCatmap();
  renderLeadmap();
}

function productCard(p){
  const caps=DATA.capOrder.map(c=>`<span class="dot ${capClass(p.capabilities[c])}" title="${DATA.capLabel[c]}: ${p.capabilities[c]}"></span> ${DATA.capLabel[c]}`).join(' &nbsp; ');
  return `<div class="card">
    <div class="rankrow">
      <div class="ranknum">${p.position}</div>
      <div style="flex:1">
        <div><span class="pos">${p.name}</span> · ${p.company} · <span class="grade g${p.evidenceGrade}">Grade ${p.evidenceGrade}</span> · ${p.confidence} confidence</div>
        <div class="note">${p.blurb}</div>
        <div style="margin:8px 0">${p.categories.map(c=>`<span class="tag cat">${c}</span>`).join('')}${p.leadership.map(l=>`<span class="tag lead">${l}</span>`).join('')}${p.surfaces.map(s=>`<span class="tag surf">${s}</span>`).join('')}</div>
        <div style="font-size:12.5px;color:var(--muted)">Capabilities — ${caps}</div>
      </div>
    </div>
  </div>`;
}

function renderCatmap(){
  const list=filtered();
  const active=new Set(list.map(p=>p.name));
  $('#catmapView').innerHTML = DATA.categories.map(cat=>{
    const members=cat.members.filter(m=>active.has(m));
    return `<div class="card"><div style="font-weight:600;margin-bottom:4px">${cat.name}</div><div class="note" style="margin-bottom:8px">${cat.desc}</div>${members.map(m=>`<span class="tag cat">${m}</span>`).join('')||'<span class="note">— (filtered out)</span>'}</div>`;
  }).join('');
}

function renderLeadmap(){
  const list=filtered();
  const active=new Set(list.map(p=>p.name));
  // group by leadership token
  const map={};
  DATA.products.forEach(p=>p.leadership.forEach(l=>{(map[l]=map[l]||[]).push(p.name);}));
  $('#leadmapView').innerHTML=Object.keys(map).sort().map(l=>{
    const members=map[l].filter(m=>active.has(m));
    return `<div class="card"><span class="tag lead">${l}</span> ${members.map(m=>`<span class="tag">${m}</span>`).join('')||'<span class="note">— (filtered out)</span>'}</div>`;
  }).join('');
}

function renderAll(){
  $('#asof').textContent='Research snapshot '+DATA.asOf+' · cutoff '+DATA.cutoff+' · rendering only (no new research)';
  $('#execJudgment').innerHTML='<strong>Executive judgment.</strong> '+DATA.execJudgment;
  $('#execStats').innerHTML=[
    ['Top 10 products','Locked selection, order fixed'],
    ['Market denominator','No single global share table'],
    ['Strategic direction','Model → Agent System → Workflow'],
    ['Capability layers','8 ordered strategic layers'],
    ['Scenario support','B > C > A (judgment, not forecast)'],
    ['Evidence registry','36 claims · 131 sources'],
  ].map(([k,v])=>`<div class="card"><div class="pos">${v}</div><div class="note">${k}</div></div>`).join('');

  $('#marketStructure').innerHTML=DATA.marketStructure.map(m=>`<div class="card"><strong>${m.k}.</strong> ${m.v}</div>`).join('');

  $('#layersTable').innerHTML='<tr><th>Layer</th><th>Importance</th><th>Differentiation</th><th>Commoditization risk</th><th>Potential moat</th><th>Confidence</th></tr>'+
    DATA.strategicLayers.map(l=>`<tr><td><strong>${l.name}</strong></td><td>${l.importance}</td><td>${l.differentiation}</td><td>${l.commoditization}</td><td>${l.moat}</td><td>${l.confidence}</td></tr>`).join('');

  $('#commodView').innerHTML=Object.keys(DATA.commoditization).map((k,i)=>`<div class="card"><div style="font-weight:600;margin-bottom:6px">${k}</div>${DATA.commoditization[k].map(x=>`<div class="tag">${x}</div>`).join('')}</div>`).join('');

  $('#compView').innerHTML=DATA.competitive.map(g=>`<div class="card"><strong>${g.group}.</strong><div style="margin-top:6px">${g.pairs.map(p=>`<span class="tag">${p}</span>`).join('')}</div></div>`).join('');

  $('#wfView').innerHTML=DATA.workflow.map((s,i)=>(i? '<span class="arrow">→</span>':'')+`<span class="s">${s}</span>`).join('');
  $('#wfNote').textContent='Already happened: code generation → repo-level execution. Emerging: parallel agent portfolios. Not yet proven: reliable end-to-end autonomous delivery (C035).';

  $('#scenView').innerHTML=DATA.scenarios.map(s=>`<div class="scenario"><strong>${s.name}</strong><div class="note" style="margin:4px 0">${s.desc}</div><div class="note">${s.note}</div></div>`).join('');
  $('#scenOrder').textContent=DATA.scenarioOrder;

  $('#riskView').innerHTML=DATA.risks.map(r=>`<div class="card"><strong>${r.name}.</strong> ${r.desc}</div>`).join('');

  $('#claimsList').innerHTML=DATA.claims.map(c=>`<div class="claimrow" data-id="${c.id}"><span class="cid">${c.id}</span><div><div>${c.claim}</div><div class="note">${c.type} · <span class="grade g${c.grade}">${c.grade}</span> · ${c.confidence} · ${c.phase} · ${c.scope}</div></div></div>`).join('');
  $$('#claimsList .claimrow').forEach(r=>r.addEventListener('click',()=>openClaim(r.dataset.id)));

  $('#foot').textContent='2026 AI Coding Agent Landscape Explorer — Phase 8 Task 4 presentation asset. Rendered from 08-dataset/products.csv, candidates.csv, 08-sources.md and 08-research-note.md. Phase 0–7 research files unmodified.';
}

function openClaim(id){
  const c=DATA.claims.find(x=>x.id===id); if(!c)return;
  $('#modalBody').innerHTML=`<span class="close" onclick="closeModal()">×</span>
    <h3>${c.id} — ${c.type}</h3>
    <div class="kv">
      <div class="k">Claim</div><div>${c.claim}</div>
      <div class="k">Evidence</div><div>${c.evidence}</div>
      <div class="k">Source</div><div><a href="../08-sources.md" style="color:var(--accent)">${c.sourceId} — ${c.source}</a></div>
      <div class="k">Date</div><div>${c.date}</div>
      <div class="k">Evidence Grade</div><div><span class="grade g${c.grade}">${c.grade}</span></div>
      <div class="k">Confidence</div><div>${c.confidence}</div>
      <div class="k">Phase</div><div>${c.phase}</div>
      <div class="k">Scope</div><div>${c.scope}</div>
    </div>`;
  $('#modal').classList.remove('hidden');
}
function closeModal(){$('#modal').classList.add('hidden');}
$('#modal').addEventListener('click',e=>{if(e.target.id==='modal')closeModal();});
document.addEventListener('keydown',e=>{if(e.key==='Escape')closeModal();});

buildFilters();
renderAll();
renderProducts();
</script>
</body>
</html>
"""

html_out = HTML_HEAD + "<script>const DATA=" + json.dumps(DATA, ensure_ascii=False) + ";</script>" + HTML_BODY
with open(OUT_HTML, "w", encoding="utf-8") as f:
    f.write(html_out)

print("Wrote", OUT_HTML, "(%d bytes)" % len(html_out))

# ---------------------------------------------------------------------------
# 5. PPTX rendering
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x12, 0x18, 0x2A)
PANEL = RGBColor(0x1B, 0x22, 0x30)
INK = RGBColor(0x12, 0x18, 0x2A)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
MUTED = RGBColor(0x5C, 0x6B, 0x80)
ACCENT = RGBColor(0x2F, 0x6F, 0xE0)
ACCENT2 = RGBColor(0x12, 0x9A, 0x7A)
WARN = RGBColor(0xB5, 0x7A, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=INK, gap=6, bold_lead=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.alignment = PP_ALIGN.LEFT
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead + "  "
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = ACCENT; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def title_bar(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Pt(3), ACCENT)
    txt(slide, Inches(0.55), Inches(0.18), SW - Inches(1.1), Inches(0.8), title,
        size=28, color=LIGHT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(slide, Inches(0.57), Inches(0.78), SW - Inches(1.1), Inches(0.3), kicker,
            size=12, color=RGBColor(0x9F, 0xB4, 0xD6))


def footer(slide, n):
    txt(slide, Inches(0.4), SH - Inches(0.42), Inches(9), Inches(0.3),
        "2026 AI Coding Agent Landscape — Executive Summary · Phase 8 Task 4 · rendering only, no new research",
        size=9, color=MUTED)
    txt(slide, SW - Inches(1.2), SH - Inches(0.42), Inches(0.8), Inches(0.3), str(n),
        size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# 1. Title
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.05), SW, Pt(3), ACCENT)
txt(s, Inches(0.9), Inches(1.9), SW - Inches(1.8), Inches(1.0),
    "2026 AI Coding Agent Landscape", size=44, color=LIGHT, bold=True)
txt(s, Inches(0.92), Inches(3.2), SW - Inches(1.8), Inches(0.6),
    "Executive Summary", size=26, color=RGBColor(0x9F, 0xB4, 0xD6))
txt(s, Inches(0.92), Inches(3.95), SW - Inches(1.8), Inches(0.5),
    "Market Structure · Product Paradigms · Agent Architecture · Workflow Evolution · Strategic Outlook",
    size=15, color=RGBColor(0xC7, 0xD3, 0xE8))
txt(s, Inches(0.92), Inches(6.4), SW - Inches(1.8), Inches(0.5),
    "Research snapshot 2026-08-31 · cutoff August 2026 · Locked Phase 3 Top 10 · not a current global ranking",
    size=12, color=MUTED)

# 2. Executive Judgment
s = add_slide(); title_bar(s, "Executive Judgment", "Core conclusion of the Case")
rect(s, Inches(0.55), Inches(1.5), SW - Inches(1.1), Inches(1.7), PANEL)
txt(s, Inches(0.8), Inches(1.65), SW - Inches(1.6), Inches(1.4),
    DATA["execJudgment"], size=18, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.6), Inches(3.5), SW - Inches(1.2), Inches(3.2), [
    ("Umbrella market", "AI Coding Agent = agentic software-engineering systems; substrate converging, product boundary diverging (C006)."),
    ("Direction", "Model → Agent System → Workflow: durable differentiation is moving above the model layer (C007)."),
    ("Market denominator", "No defensible single global market-share table for Aug 2026; JetBrains figures are multi-select, not shares (C002, C005)."),
    ("Productivity", "Benchmarks support capability, not productivity; automated pass overstates maintainer acceptance ~24.2 pts (C013, C014)."),
], size=15)
footer(s, 2)

# 3. Market Structure
s = add_slide(); title_bar(s, "Market Structure", "Umbrella market — strata are not quality tiers")
y = Inches(1.5)
for m in DATA["marketStructure"]:
    rect(s, Inches(0.55), y, SW - Inches(1.1), Inches(1.55), PANEL)
    txt(s, Inches(0.8), y + Inches(0.1), SW - Inches(1.6), Inches(0.4), m["k"], size=16, color=ACCENT, bold=True)
    txt(s, Inches(0.8), y + Inches(0.55), SW - Inches(1.6), Inches(0.9), m["v"], size=13.5, color=INK)
    y += Inches(1.75)
footer(s, 3)

# 4. Top 10
s = add_slide(); title_bar(s, "Top 10", "Locked Phase 3 selection — Market × Technology/Product Significance")
txt(s, Inches(0.55), Inches(1.22), SW - Inches(1.1), Inches(0.3),
    "Order is locked; do not re-sort or present as a current global ranking.", size=12, color=WARN, italic=True)
cols = 2; rows = 5; x0 = Inches(0.55); y0 = Inches(1.6)
cw = int((SW - Inches(1.1) - Inches(0.3)) / cols); ch = Inches(1.02)
for i, p in enumerate(DATA["products"]):
    c = i % cols; r = i // cols
    x = x0 + c * (cw + Inches(0.3)); y = y0 + r * (ch + Inches(0.08))
    rect(s, x, y, cw, ch, PANEL)
    txt(s, x + Inches(0.12), y + Inches(0.05), cw - Inches(0.24), Inches(0.4),
        f"{p['position']}.  {p['name']}", size=16, color=ACCENT, bold=True)
    txt(s, x + Inches(0.12), y + Inches(0.42), cw - Inches(0.24), Inches(0.3),
        p["company"], size=11, color=MUTED)
    txt(s, x + Inches(0.12), y + Inches(0.68), cw - Inches(0.24), Inches(0.3),
        p["paradigm"], size=11.5, color=INK)
footer(s, 4)

# 5. Product Paradigms
s = add_slide(); title_bar(s, "Product Paradigms", "Eight overlapping architectures (Phase 6)")
y = Inches(1.45)
for pr in DATA["paradigms"]:
    txt(s, Inches(0.6), y, Inches(3.7), Inches(0.6), pr["name"], size=14.5, color=ACCENT, bold=True)
    txt(s, Inches(4.4), y, SW - Inches(4.9), Inches(0.6), pr["desc"], size=12.5, color=INK)
    y += Inches(0.68)
footer(s, 5)

# 6. Leadership Map
s = add_slide(); title_bar(s, "Leadership Map", "Category judgments — not a second overall ranking")
lead_map = {}
for p in DATA["products"]:
    for l in p["leadership"]:
        lead_map.setdefault(l, []).append(p["name"])
y = Inches(1.5)
for l in sorted(lead_map):
    txt(s, Inches(0.6), y, Inches(5.2), Inches(0.5), l, size=14, color=ACCENT, bold=True)
    txt(s, Inches(6.0), y, SW - Inches(6.4), Inches(0.5), " · ".join(lead_map[l]), size=13.5, color=INK)
    y += Inches(0.6)
footer(s, 6)

# 7. Strategic Layer Model
s = add_slide(); title_bar(s, "Strategic Layer Model", "8 ordered layers (Phase 7) — ordered, not ranked")
rows = [["Layer", "Importance", "Differentiation", "Commoditization", "Potential moat"]]
for l in DATA["strategicLayers"]:
    rows.append([l["name"], l["importance"], l["differentiation"], l["commoditization"], l["moat"]])
nrows = len(rows); ncols = 5
tbl_x = Inches(0.55); tbl_y = Inches(1.5); tbl_w = SW - Inches(1.1); tbl_h = Inches(5.4)
gtbl = s.shapes.add_table(nrows, ncols, tbl_x, tbl_y, tbl_w, tbl_h).table
gtbl.columns[0].width = Inches(2.2)
for ci in range(1, 5):
    gtbl.columns[ci].width = int((tbl_w - Inches(2.2)) / 4)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.text = val
        para = cell.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        run = para.runs[0]; run.font.size = Pt(12 if ri == 0 else 11)
        run.font.bold = (ri == 0 or ci == 0)
        run.font.name = "Calibri"
        run.font.color.rgb = LIGHT if ri == 0 else INK
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if ri == 0 else (PANEL if ri % 2 else RGBColor(0x21, 0x29, 0x39))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(5); cell.margin_right = Pt(5); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
footer(s, 7)

# 8. Capability Commoditization
s = add_slide(); title_bar(s, "Capability Commoditization", "3-level structure (Phase 7) — no added capabilities")
labels = list(DATA["commoditization"].keys())
colw = int((SW - Inches(1.1) - Inches(0.6)) / 3)
for i, k in enumerate(labels):
    x = Inches(0.55) + i * (colw + Inches(0.3))
    rect(s, x, Inches(1.5), colw, Inches(0.55), ACCENT if i == 2 else NAVY)
    txt(s, x, Inches(1.5), colw, Inches(0.55), k, size=13.5, color=LIGHT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    items = DATA["commoditization"][k]
    bullets(s, x + Inches(0.05), Inches(2.2), colw - Inches(0.1), Inches(4.6), items, size=12.5, gap=7)
footer(s, 8)

# 9. Competitive Structure
s = add_slide(); title_bar(s, "Competitive Structure", "Relationships from Phase 6 only — typed, not ranked")
y = Inches(1.5)
for g in DATA["competitive"]:
    txt(s, Inches(0.6), y, Inches(3.0), Inches(0.5), g["group"], size=15, color=ACCENT, bold=True)
    txt(s, Inches(3.7), y, SW - Inches(4.2), Inches(0.5), "   ".join(g["pairs"]), size=13, color=INK)
    y += Inches(0.7)
txt(s, Inches(0.6), y + Inches(0.1), SW - Inches(1.2), Inches(0.5),
    "No new competitive relationships invented.", size=12, color=WARN, italic=True)
footer(s, 9)

# 10. Workflow Evolution
s = add_slide(); title_bar(s, "Workflow Evolution", "Work unit rising — not yet whole-workflow autonomy")
chain = "  →  ".join(DATA["workflow"])
rect(s, Inches(0.55), Inches(1.6), SW - Inches(1.1), Inches(1.2), PANEL)
txt(s, Inches(0.8), Inches(1.6), SW - Inches(1.6), Inches(1.2), chain, size=15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.6), Inches(3.2), SW - Inches(1.2), Inches(3.0), [
    ("Already happened", "code generation → repo-level execution; terminal/runtime interaction standard; background agents productized."),
    ("Emerging now", "human supervision of parallel agent portfolios; runtime provisioning as part of autonomy; persistent memory as context."),
    ("Not yet proven", "reliable low-supervision end-to-end autonomous delivery (C035)."),
], size=15)
footer(s, 10)

# 11. Strategic Scenarios
s = add_slide(); title_bar(s, "Strategic Scenarios", "Phase 7 scenarios only — no forecasts")
y = Inches(1.5)
for sc in DATA["scenarios"]:
    rect(s, Inches(0.55), y, SW - Inches(1.1), Inches(1.45), PANEL)
    txt(s, Inches(0.8), y + Inches(0.08), SW - Inches(1.6), Inches(0.4), sc["name"], size=15, color=ACCENT, bold=True)
    txt(s, Inches(0.8), y + Inches(0.5), SW - Inches(1.6), Inches(0.6), sc["desc"], size=12.5, color=INK)
    txt(s, Inches(0.8), y + Inches(1.06), SW - Inches(1.6), Inches(0.3), sc["note"], size=11, color=MUTED, italic=True)
    y += Inches(1.6)
txt(s, Inches(0.6), y - Inches(0.05), SW - Inches(1.2), Inches(0.4), DATA["scenarioOrder"], size=12.5, color=WARN, bold=True)
footer(s, 11)

# 12. Risks / Unknowns
s = add_slide(); title_bar(s, "Risks / Unknowns", "Evidence limits carried forward, never resolved")
half = int((SW - Inches(1.1) - Inches(0.4)) / 2)
for i, r in enumerate(DATA["risks"]):
    col = i % 2; row = i // 2
    x = Inches(0.55) + col * (half + Inches(0.4)); y = Inches(1.5) + row * Inches(1.75)
    rect(s, x, y, half, Inches(1.6), PANEL)
    txt(s, x + Inches(0.12), y + Inches(0.08), half - Inches(0.24), Inches(0.4), r["name"], size=13.5, color=ACCENT, bold=True)
    txt(s, x + Inches(0.12), y + Inches(0.5), half - Inches(0.24), Inches(1.0), r["desc"], size=11, color=INK)
footer(s, 12)

# 13. Final Takeaway
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.5), SW, Pt(3), ACCENT)
txt(s, Inches(0.9), Inches(0.9), SW - Inches(1.8), Inches(0.7), "Final Takeaway", size=34, color=LIGHT, bold=True)
bullets(s, Inches(0.95), Inches(2.8), SW - Inches(1.9), Inches(3.6), [
    ("Umbrella market", "substrate converging, product boundary diverging — compare by paradigm, not feature count."),
    ("Differentiation", "moving to harness, runtime, context/memory, orchestration, verification and workflow integration."),
    ("Top 10 is locked", "a research selection, not a live leaderboard; vendor scale figures stay labeled as vendor claims."),
    ("Honest limits", "no global share table, autonomy not proven, productivity not cleanly measured — these are standing unknowns."),
], size=17, color=LIGHT, gap=12)
txt(s, Inches(0.95), Inches(6.7), SW - Inches(1.9), Inches(0.5),
    "Source: 08-dataset (products.csv, candidates.csv), 08-sources.md, 08-research-note.md. Phase 0–7 unmodified.",
    size=11, color=MUTED)

prs.save(OUT_PPTX)
print("Wrote", OUT_PPTX, "(%d slides)" % len(prs.slides._sldIdLst))
